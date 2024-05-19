import streamlit as st
from pptx import Presentation
from docx import Document
from transformers import pipeline
import fitz
import time

# Load pretrained summarization model
summarizer = pipeline("summarization")

def extract_text_from_ppt(ppt):
    pr = Presentation(ppt)
    text = []
    for slide in pr.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_docx(file):
    text = ""
    doc = Document(file)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pdf(uploaded_file):
    text = ""
    pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        text += page.get_text()
    return text

fun_facts = [
    "Fools who don’t respect the past are likely to repeat it.   - Nico Robin",
    "No matter how hard or impossible it is, never lose sight of your goal.  - Monkey D Luffy.",
    "The difference between the novice and the master is that the master has failed more times than the novice has tried.  -Koro-sensei",
    "I never treated anyone as ally. they are just tools. All that matters in winning- Ayanokyoji Kiyotaka",
    "The world is curel but also beautiful. - Mikasa Ackerman "
]

st.set_page_config(page_title="Document Text Summarizer", page_icon="🗒️", layout="wide")
st.title("🗒️  SUMMARY GENERATOR CHAN")

# Creating a two-column layout
col1, col2 = st.columns([3, 1])

with col1:
    try:
        st.image('study.jpg', caption='THE ONE AND ONLY')
    except FileNotFoundError:
        st.warning("Image file 'study.jpg' not found.")

    # File uploader for PPTX, DOCX, or PDF files
    uploaded_file = st.file_uploader("Upload a PPT, DOCX, or PDF document", type=["pptx", "docx", "pdf"])

    # Text input for direct text entry
    user_input_text = st.text_area("Input Text for Summarization", "")

    if st.button("Generate Summary"):
        if uploaded_file:
            file_ext = uploaded_file.name.split(".")[-1]
            if file_ext == "pptx":
                extracted_text = extract_text_from_ppt(uploaded_file)
            elif file_ext == "docx":
                extracted_text = extract_text_from_docx(uploaded_file)
            elif file_ext == "pdf":
                extracted_text = extract_text_from_pdf(uploaded_file)
        elif user_input_text:
            extracted_text = user_input_text
        else:
            st.write("Please upload a file or enter some text for summarization.")

        if extracted_text:
            st.write("Enjoy the music while you wait...")
            # Play audio
            st.audio('audio.mp3', format='audio/mp3', start_time=0)
            
            # Generate summary
            with st.spinner('Generating summary...'):
                time.sleep(5)  # Simulate time delay
                summary = summarizer(extracted_text, max_length=500, min_length=100, do_sample=False)
                # Adjust summary to ensure it falls within the desired range
                if len(summary[0]['summary_text'].split()) < 100:
                    st.warning("Summary is too short. Increasing length...")
                    summary = summarizer(extracted_text, max_length=500, min_length=100, do_sample=False)
                elif len(summary[0]['summary_text'].split()) > 500:
                    st.warning("Summary is too long. Decreasing length...")
                    summary = summarizer(extracted_text, max_length=500, min_length=100, do_sample=False)
                st.write("SUMMARY :")
                st.write(summary[0]['summary_text'])

# Fun facts in the second column
with col2:
    st.write("## Anime Quotes")
    for fact in fun_facts:
        st.markdown(f"""
        <div style="background-color: #ffeb3b; border-radius: 10px; padding: 10px; margin-bottom: 10px; color: black;">
            {fact}
        </div>
        """, unsafe_allow_html=True)
